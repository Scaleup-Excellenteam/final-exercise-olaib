from flask import json
from pptx import Presentation


class PptxProcessor:
    """Parser class for extracting text from PowerPoint(pptx) slides."""

    def __init__(self, path: str):
        self.path = path
        self.file_name = path.split("/")[-1].split(".")[0]
        self.slides = self.extract_text_from_presentation()

    def extract_text_from_presentation(self):
        """Extract text from all slides in the presentation.
        :return: list of slide texts
        """
        presentation = Presentation(self.path)
        slides_text = []
        for slide in presentation.slides:
            slide_text = self.extract_text_from_slide(slide)
            slides_text.append(slide_text)
        return slides_text

    def get_slide(self, slide_num) -> str:
        return self.slides[slide_num - 1]

    def get_slides(self) -> list:
        return self.slides

    @staticmethod
    def extract_text_from_slide(slide):
        """Extract text from a single slide.
        :param slide: slide object
        """
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text += run.text
                        slide_text += " ."  # Space between paragraphs
        return slide_text

    def save_explanations_as_json(self, explanations):
        """Save explanations as a json file.
        :param explanations: list of explanations
        """
        data = {}
        for i, explanation in enumerate(explanations):
            slide_number = i + 1
            data[slide_number] = explanation

        with open(f"explanations_{self.file_name}.json", "w") as file:
            json.dump(data, file, indent=4)

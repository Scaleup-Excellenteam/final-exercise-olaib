import json
import openai
from pptx import Presentation
import os

# API_KEY = os.environ.get('API_KEY')
# MODEL_ID = os.environ.get('MODEL_ID')

API_KEY = "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"
MODEL_ID = "gpt-3.5-turbo"


class SingletonMeta(type):
    _instances = {}

    def __call__(cls, *args, **kwargs):
        if cls not in cls._instances:
            cls._instances[cls] = super().__call__(*args, **kwargs)
        return cls._instances[cls]


class OpenAIProcessor(metaclass=SingletonMeta):
    def __init__(self):
        self.api_key = API_KEY
        self.model_id = MODEL_ID
        self.initialize_openai()

    def initialize_openai(self):
        openai.api_key = self.api_key

    def generate_explanation(self, prompt):
        print(prompt)
        response = openai.ChatCompletion.create(
            model=self.model_id,
            messages=[
                {"role": "system", "content": "You are a presentation assistant."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=100,
            temperature=0.9,
            top_p=1,
            frequency_penalty=0.0,
            presence_penalty=0.6,
            stop=["\n", " Slide Text:"]
        )
        print(response)
        # extract answer from response
        if response and response.choices:
            return response.choices[0].message.get('content', '')
        return ''


def extract_text_from_slide(slide):
    slide_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text += run.text
                    slide_text += " "  # Add space between paragraphs
    return slide_text


class Presentation_processor:
    def __init__(self):
        self.presentation = None
        self.slides = []
        self.explanations = []
        self.openai_processor = OpenAIProcessor()

    def open_presentation(self, path: str):
        self.presentation = Presentation(path)

    def generate_explanation(self, slide_text):
        slide_text = f"Slide Text: {slide_text}"
        explanation = self.openai_processor.generate_explanation(slide_text)
        return explanation

    def process_presentation(self):
        for slide_num, slide in enumerate(self.presentation.slides):
            slide_text = extract_text_from_slide(slide)
            explanation = self.generate_explanation(slide_text)
            self.slides.append(slide_text)
            self.explanations.append(explanation)

    def save_explanations(self, output_file):
        data = {}
        for i, explanation in enumerate(self.explanations):
            slide_number = i + 1
            data[slide_number] = explanation

        with open(output_file, "w") as file:
            json.dump(data, file, indent=4)

        print(f"Explanations saved successfully to {output_file}")

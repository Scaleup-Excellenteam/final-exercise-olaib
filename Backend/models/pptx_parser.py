import os

from flask import json
from pptx import Presentation

from config import app_log as logger


def save_content_to_json_file(contents, output_folder, file_name: str) -> None:
    data = {}
    for i, content in enumerate(contents):
        slide_number = i + 1
        data[slide_number] = content
    logger.info(f"data as json: {data}")
    # Save as json file into output folder
    filename = file_name + ".json"
    with open(os.path.join(output_folder, filename), 'w') as outfile:
        json.dump(data, outfile)


class PptxProcessor:
    """Parser class for extracting text from PowerPoint(pptx) slides."""

    # # get file as pptx and filename as str
    # def __init__(self, file, file_name: str):
    #     self.file = file
    #     self.file_name = file_name
    #     self.slides = self.extract_text_from_presentation()
    #
    # def save_parse_content_to_json(self, output_file):
    #     self.save_content_to_json_file(self.get_slides(), output_file)
    #
    # def extract_text_from_presentation(self):
    #     """Extract text from all slides in the presentation.
    #     :return: list of slide texts
    #     """
    #     presentation = Presentation(self.file)
    #     slides_text = []
    #     for slide in presentation.slides:
    #         slide_text = self.extract_text_from_slide(slide)
    #         slides_text.append(slide_text)
    #     return slides_text
    #
    # def get_slide(self, slide_num) -> str:
    #     return self.slides[slide_num - 1]
    #
    # def get_slides(self) -> list:
    #     return self.slides
    #
    # @staticmethod
    # def extract_text_from_slide(slide):
    #     """Extract text from a single slide.
    #     :param slide: slide object
    #     """
    #     slide_text = ""
    #     for shape in slide.shapes:
    #         if shape.has_text_frame:
    #             for paragraph in shape.text_frame.paragraphs:
    #                 for run in paragraph.runs:
    #                     slide_text += run.text
    #                     slide_text += " ."  # Space between paragraphs
    #     return slide_text
    #
    # def save_explanations_as_json(self, output_folder):
    #     """Save explanations as a JSON file.
    #     :param output_folder: output folder to save the file
    #     :param self: list of explanations
    #     """
    #     self.save_content_to_json_file(self, output_folder)
    #
    # @staticmethod
    # def save_content_to_json_file(contents, upload_folder):
    #     data = {}
    #     for i, content in enumerate(contents):
    #         slide_number = i + 1
    #         data[slide_number] = content
    #     logger.info(f"data as json: {data}")
    #
    #     # Save as json file into output folder
    #     # save as json file ,if file exist overwrite it with new data, if not create new file
    #     with open(upload_folder, 'w') as outfile:
    #         json.dump(data, outfile)
    """Parser class for extracting text from PowerPoint(pptx) slides."""

    # get file as pptx and filename as str
    def __init__(self, file, file_name: str):
        self.file = file
        self.file_name = file_name
        self.slides = self.extract_text_from_presentation()

    def save_parse_content_to_json(self, upload_file):
        save_content_to_json_file(self.get_slides(), upload_file, self.file_name)

    def extract_text_from_presentation(self):
        """Extract text from all slides in the presentation.
        :return: list of slide texts
        """
        presentation = Presentation(self.file)
        slides_text = []
        for slide in presentation.slides:
            slide_text = self.extract_text_from_slide(slide)
            slides_text.append(slide_text)
        return slides_text

    def get_slide(self, slide_num) -> str:
        return self.slides[slide_num - 1]

    def get_slides(self) -> list:
        return self.slides

    @staticmethod
    def extract_text_from_slide(slide):
        """Extract text from a single slide.
        :param slide: slide object
        """
        slide_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text += run.text
                        slide_text += " ."  # Space between paragraphs
        return slide_text

    def save_explanations_as_json(self, explanations, output_folder) -> None:
        """Save explanations as a JSON file.
        :param output_folder: output folder to save the file
        :param explanations: list of explanations
        """
        save_content_to_json_file(explanations, output_folder, self.file_name)
